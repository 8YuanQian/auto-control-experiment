# -*- coding: utf-8 -*-
'''
四旋翼无人机姿态稳定控制系统设计 —— PPT 汇报
基于 MATLAB/Simulink 仿真，包含模型截图、代码截图、结果对比图
'''

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 路径配置 ====================
DESKTOP = os.path.expanduser('~/Desktop')
SCREENSHOT_DIR = os.path.join(DESKTOP, 'professional lesson', 'matlab_sim', 'screenshots')
OUTPUT = os.path.join(DESKTOP, '四旋翼无人机姿态控制_课设汇报_Simulink.pptx')

# ==================== 配色方案 (暗金学术) ====================
BG_DARK = RGBColor(0x0A, 0x0F, 0x14)
BG_SLIDE = RGBColor(0x16, 0x21, 0x2B)
GOLD = RGBColor(0xD4, 0xA0, 0x1E)
GOLD_LIGHT = RGBColor(0xF0, 0xC8, 0x40)
WHITE = RGBColor(0xE8, 0xEC, 0xF0)
WHITE_DIM = RGBColor(0xA0, 0xA8, 0xB4)
ACCENT_BLUE = RGBColor(0x40, 0x90, 0xD0)
ACCENT_RED = RGBColor(0xE0, 0x50, 0x50)
ACCENT_GREEN = RGBColor(0x40, 0xC0, 0x70)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_gold_line(slide, left, top, width):
    '''添加金色装饰线'''
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_gold_bar_left(slide):
    '''左侧金色竖条'''
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.08), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def add_pic(slide, filename, left, top, width, height=None):
    '''添加图片'''
    path = os.path.join(SCREENSHOT_DIR, filename)
    if not os.path.exists(path):
        print(f'  [WARN] 图片不存在: {path}')
        return None
    if height:
        return slide.shapes.add_picture(path, Inches(left), Inches(top),
                                        Inches(width), Inches(height))
    else:
        return slide.shapes.add_picture(path, Inches(left), Inches(top),
                                        Inches(width))


def add_page_number(slide, num, total):
    add_textbox(slide, 12.0, 7.05, 1.2, 0.35,
                f'{num}/{total}', font_size=10, color=WHITE_DIM,
                alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12

# =====================================================
# Slide 1: 封面
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_gold_bar_left(slide)
add_gold_line(slide, 0.5, 3.55, 12.3)

add_textbox(slide, 0.8, 1.0, 11.5, 0.7,
            '四旋翼无人机姿态稳定控制系统设计', font_size=42, color=GOLD, bold=True)
add_textbox(slide, 0.8, 1.7, 11.5, 0.6,
            '基于 Simulink 的建模、PID/LQR 设计与仿真', font_size=24, color=WHITE)

add_textbox(slide, 0.8, 4.2, 11.5, 0.5,
            '自动控制原理 · 课程设计汇报', font_size=22, color=WHITE_DIM)
add_textbox(slide, 0.8, 5.0, 5.0, 0.4,
            '学生：陈培文  202400405059', font_size=18, color=WHITE_DIM)
add_textbox(slide, 0.8, 5.5, 5.0, 0.4,
            '任课教师：程涛  郭晓东', font_size=18, color=WHITE_DIM)
add_textbox(slide, 0.8, 6.0, 5.0, 0.4,
            '深圳技术大学 · 2026.06.14', font_size=18, color=WHITE_DIM)

# =====================================================
# Slide 2: 目录
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)

add_textbox(slide, 0.8, 0.3, 5.0, 0.6, '目 录', font_size=36, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

toc_items = [
    ('01', '研究背景与意义', '四旋翼应用场景、姿态控制核心地位'),
    ('02', '问题描述与动力学建模', 'Newton-Euler 建模 -> 传递函数 Gφ(s)'),
    ('03', 'PID 控制器设计', '根轨迹法整定、Simulink 模型搭建'),
    ('04', 'LQR 控制器设计', 'Riccati 方程求解、全状态反馈'),
    ('05', 'Simulink 仿真对比', '阶跃响应 / 抗扰动 / 参数摄动'),
    ('06', '创新改进', '前馈 / 增益调度 / 抗积分饱和 / 串级'),
    ('07', '总结与展望', '结论 + 未来方向'),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = 1.5 + i * 0.75
    add_textbox(slide, 0.8, y, 0.8, 0.4, num, font_size=28, color=GOLD, bold=True)
    add_textbox(slide, 1.7, y, 4.0, 0.4, title, font_size=20, color=WHITE, bold=True)
    add_textbox(slide, 1.7, y + 0.35, 6.0, 0.3, desc, font_size=13, color=WHITE_DIM)

# 右侧放置一张模型总览图
add_pic(slide, '04_pid_simulink_model.png', 7.5, 1.2, 5.0, 5.5)

# =====================================================
# Slide 3: 研究背景与意义
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)
add_textbox(slide, 0.8, 0.3, 10.0, 0.6, '01  研究背景与意义', font_size=32, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

items = [
    '四旋翼无人机在航拍、植保、物流、应急救援等领域广泛应用',
    '姿态稳定是飞控最底层环节——姿态环出问题，速度/位置环无从谈起',
    '四旋翼姿态系统：MIMO、非线性、强耦合、欠驱动',
    '核心挑战：电机延迟（τ ≈ 0.05–0.15 s）不可忽略',
    '目标：将建模-分析-设计-验证闭环走通，覆盖自控原理核心知识点',
]
for i, item in enumerate(items):
    add_textbox(slide, 0.8, 1.5 + i * 0.9, 5.5, 0.7, f'● {item}',
                font_size=17, color=WHITE)

add_pic(slide, '01_open_loop_analysis.png', 7.0, 1.2, 5.5, 5.5)

# =====================================================
# Slide 4: 动力学建模
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)
add_textbox(slide, 0.8, 0.3, 12.0, 0.6, '02  问题描述与动力学建模', font_size=32, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

add_textbox(slide, 0.8, 1.3, 6.0, 0.4, '■ 物理参数（DJI F450 + 2212电机）', font_size=18, color=GOLD, bold=True)

params = [
    '总质量 m = 1.0 kg，机臂长度 l = 0.225 m',
    '绕 x 轴转动惯量 Ixx = 0.012 kg·m²',
    '电机增益 Km = 1.0 N/V，时间常数 τ = 0.1 s',
    '滚转通道开环增益 K = Km·l / Ixx = 18.75',
]
for i, p in enumerate(params):
    add_textbox(slide, 0.8, 1.8 + i * 0.45, 6.0, 0.35, f'  {p}', font_size=15, color=WHITE)

add_textbox(slide, 0.8, 3.8, 6.0, 0.4, '■ 滚转通道传递函数', font_size=18, color=GOLD, bold=True)
add_textbox(slide, 0.8, 4.3, 6.0, 0.4, '  Gφ(s) = 18.75 / [s²·(0.1s + 1)]', font_size=18,
            color=ACCENT_BLUE, bold=True)

add_textbox(slide, 0.8, 5.0, 6.0, 0.4, '■ 极点分布', font_size=18, color=GOLD, bold=True)
add_textbox(slide, 0.8, 5.5, 6.0, 0.4,
            '  s₁,₂ = 0 (原点双积分)  ->  临界稳定', font_size=15, color=WHITE)
add_textbox(slide, 0.8, 5.9, 6.0, 0.4,
            '  s₃ = −10 (电机惯性极点)', font_size=15, color=WHITE)

add_pic(slide, '01_open_loop_analysis.png', 6.8, 1.2, 5.8, 5.8)

# =====================================================
# Slide 5: PID 控制器设计
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)
add_textbox(slide, 0.8, 0.3, 12.0, 0.6, '03  PID 控制器设计 —— 根轨迹法', font_size=32, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

add_textbox(slide, 0.8, 1.2, 6.0, 0.4, '■ 设计指标', font_size=18, color=GOLD, bold=True)
add_textbox(slide, 0.8, 1.7, 6.0, 0.35, '  阻尼比 ζ = 0.7，自然频率 ωn = 5 rad/s', font_size=15, color=WHITE)
add_textbox(slide, 0.8, 2.1, 6.0, 0.35, '  期望主导极点: s₁,₂ = −3.5 ± j3.57', font_size=15, color=WHITE)

add_textbox(slide, 0.8, 2.7, 6.0, 0.4, '■ PID 参数', font_size=18, color=GOLD, bold=True)
add_textbox(slide, 0.8, 3.2, 6.0, 0.35, '  Kp = 0.923    Ki = 0.417    Kd = 0.187', font_size=16, color=ACCENT_BLUE, bold=True)
add_textbox(slide, 0.8, 3.6, 6.0, 0.35, '  微分滤波器 N = 100，Anti-windup ±2.0', font_size=15, color=WHITE)

add_textbox(slide, 0.8, 4.2, 6.0, 0.4, '■ 特点', font_size=18, color=GOLD, bold=True)
add_textbox(slide, 0.8, 4.7, 6.0, 0.5, '  根轨迹法整定直观，工程易复现，参数物理意义明确', font_size=15, color=WHITE)

# PID 设计分析图
add_pic(slide, '02_pid_design_analysis.png', 0.5, 5.6, 6.0, 1.8)

# PID Simulink 模型截图
add_pic(slide, '04_pid_simulink_model.png', 7.0, 1.2, 5.5, 5.0)

# =====================================================
# Slide 6: PID MATLAB 代码
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)
add_textbox(slide, 0.8, 0.3, 12.0, 0.6, '03  PID 控制器 —— MATLAB 代码与 Simulink 实现',
            font_size=32, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

add_pic(slide, '10_code_pid_design.png', 0.5, 1.3, 6.2, 4.2)
add_pic(slide, '04_pid_simulink_model.png', 7.0, 1.3, 5.8, 4.2)

add_textbox(slide, 0.8, 6.0, 12.0, 0.5,
            'Simulink 求解器: ode45 | 最大步长 0.01 s | 相对容差 1×10⁻⁶ | Stop Time 10 s',
            font_size=14, color=WHITE_DIM)

# =====================================================
# Slide 7: LQR 控制器设计
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)
add_textbox(slide, 0.8, 0.3, 12.0, 0.6, '04  LQR 控制器设计 —— 最优状态反馈', font_size=32, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

add_textbox(slide, 0.8, 1.2, 6.0, 0.4, '■ 状态空间模型', font_size=18, color=GOLD, bold=True)
add_textbox(slide, 0.8, 1.7, 6.0, 0.35, '  状态 x = [φ, ̇φ, ωm]ᵀ（角度/角速率/电机状态）', font_size=15, color=WHITE)
add_textbox(slide, 0.8, 2.1, 6.0, 0.35, '  A = [0,1,0; 0,0,18.75; 0,0,−10],  B = [0;0;10]', font_size=15, color=WHITE)

add_textbox(slide, 0.8, 2.7, 6.0, 0.4, '■ 性能指标与求解', font_size=18, color=GOLD, bold=True)
add_textbox(slide, 0.8, 3.2, 6.0, 0.35, '  Q = diag(100, 1, 1)，R = 1（重点罚角度偏差）', font_size=15, color=WHITE)
add_textbox(slide, 0.8, 3.6, 6.0, 0.35, '  lqr() 求解 Riccati 方程', font_size=15, color=WHITE)
add_textbox(slide, 0.8, 4.0, 6.0, 0.35, '  K_LQR = [10.00, 4.472, 0.287]', font_size=16, color=ACCENT_BLUE, bold=True)

add_textbox(slide, 0.8, 4.6, 6.0, 0.4, '■ 特点', font_size=18, color=GOLD, bold=True)
add_textbox(slide, 0.8, 5.1, 6.0, 0.5, '  多目标最优平衡，需全状态反馈（实际部署需 Kalman）', font_size=15, color=WHITE)

add_pic(slide, '05_lqr_simulink_model.png', 7.0, 1.2, 5.5, 5.0)
add_pic(slide, '03_lqr_design_analysis.png', 0.5, 5.8, 6.0, 1.5)

# =====================================================
# Slide 8: LQR MATLAB 代码
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)
add_textbox(slide, 0.8, 0.3, 12.0, 0.6, '04  LQR 控制器 —— MATLAB 代码与 Simulink 实现',
            font_size=32, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

add_pic(slide, '11_code_lqr_design.png', 0.5, 1.3, 6.0, 4.0)
add_pic(slide, '05_lqr_simulink_model.png', 7.0, 1.3, 5.8, 4.0)

add_textbox(slide, 0.8, 5.8, 12.0, 0.5,
            'C = eye(3) 输出全状态 [φ, ̇φ, ωm] -> Demux 分离 -> LQR Gain 矩阵反馈',
            font_size=14, color=WHITE_DIM)

# =====================================================
# Slide 9: 仿真结果对比
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)
add_textbox(slide, 0.8, 0.3, 12.0, 0.6, '05  Simulink 仿真对比 —— PID vs LQR', font_size=32, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

add_pic(slide, '07_pid_lqr_comparison.png', 0.5, 1.2, 12.3, 6.0)

# =====================================================
# Slide 10: 抗扰动对比
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)
add_textbox(slide, 0.8, 0.3, 12.0, 0.6, '05  抗扰动性能对比 & 串级控制',
            font_size=32, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

add_pic(slide, '08_disturbance_comparison.png', 0.3, 1.2, 6.4, 4.5)
add_pic(slide, '09_cascade_analysis.png', 6.8, 1.2, 6.2, 3.8)

add_pic(slide, '06_cascade_simulink_model.png', 6.8, 5.2, 5.5, 2.0)

# =====================================================
# Slide 11: 创新改进
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_gold_bar_left(slide)
add_textbox(slide, 0.8, 0.3, 12.0, 0.6, '06  创新改进', font_size=32, color=GOLD, bold=True)
add_gold_line(slide, 0.8, 0.95, 3.0)

improvements = [
    ('前馈补偿', '大角度机动时叠加 uff = Ixx·φ̈_d，跟踪误差降低 ~40%'),
    ('增益调度', '预置空载/半载/满载三组 LQR 增益，飞控根据质量估计在线切换'),
    ('抗积分饱和', '|e| > 0.1 rad 时冻结积分器，大角度阶跃超调从 14.2% -> 8.7%'),
    ('双回路串级', '角速率内环 PD（高带宽，给阻尼）+ 角度外环 P（做跟踪）'),
]
for i, (title, desc) in enumerate(improvements):
    y = 1.5 + i * 1.4
    add_textbox(slide, 0.8, y, 3.0, 0.4, f'■ {title}', font_size=20, color=GOLD, bold=True)
    add_textbox(slide, 0.8, y + 0.45, 12.0, 0.4, desc, font_size=16, color=WHITE)

# =====================================================
# Slide 12: 总结与感谢
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_gold_bar_left(slide)
add_gold_line(slide, 0.8, 3.5, 11.5)

add_textbox(slide, 0.8, 0.8, 11.5, 0.6, '07  总结与展望', font_size=32, color=GOLD, bold=True)

conclusions = [
    '① 基于 Newton-Euler 方程建立滚转通道三阶模型 Gφ(s) = 18.75/[s²(0.1s+1)]',
    '② PID (根轨迹): tr=0.82s, σ=6.3%, ess=0 —— 经典方法，工程直观',
    '③ LQR (最优): σ=2.1%, trec=0.91s, 能耗 ~60% PID —— 多目标平衡更优',
    '④ 双回路串级架构进一步改善抗扰性能，模块化建模便于迭代',
]
for i, c in enumerate(conclusions):
    add_textbox(slide, 0.8, 1.5 + i * 0.65, 11.5, 0.5, c, font_size=18, color=WHITE)

add_textbox(slide, 0.8, 4.3, 11.5, 0.5,
            '展望: 三轴联合控制 + ESO 自抗扰 + PX4 HIL 实飞验证',
            font_size=18, color=GOLD_LIGHT, bold=True)

add_textbox(slide, 0.8, 5.5, 11.5, 0.6, '感谢各位老师！', font_size=36, color=GOLD, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 6.2, 11.5, 0.4, '深圳技术大学 · 2026.06.14',
            font_size=18, color=WHITE_DIM, alignment=PP_ALIGN.CENTER)

# ==================== 保存 ====================
prs.save(OUTPUT)
print(f'\nPPT 已保存: {OUTPUT}')
print(f'幻灯片: {len(prs.slides)} 页')
